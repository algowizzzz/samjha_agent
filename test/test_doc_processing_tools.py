"""Tests for document processing MCP tools."""

import base64
import tempfile
import unittest
from pathlib import Path

try:
    import docx  # type: ignore
except ImportError:  # pragma: no cover
    docx = None

try:
    from pptx import Presentation  # type: ignore
except ImportError:  # pragma: no cover
    Presentation = None

try:
    import pdfplumber  # type: ignore
except ImportError:  # pragma: no cover
    pdfplumber = None

from external.tools.doc_processing.convert_to_markdown import ConvertToMarkdownTool
from external.tools.doc_processing.extract_images import ExtractImagesTool


_SAMPLE_PDF_BASE64 = (
    "JVBERi0xLjUKJeLjz9MKMSAwIG9iago8PAovVHlwZSAvQ2F0YWxvZwovUGFnZXMgMiAwIFIKPj4K"
    "ZW5kb2JqCjIgMCBvYmoKPDwKL1R5cGUgL1BhZ2VzCi9LaWRzIFszIDAgUl0KL0NvdW50IDEKPj4K"
    "ZW5kb2JqCjMgMCBvYmoKPDwKL1R5cGUgL1BhZ2UKL1BhcmVudCAyIDAgUgovTWVkaWFCb3ggWzAg"
    "MCA1OTUuMjggODQxLjg5XQovQ29udGVudHMgNCAwIFIKL1Jlc291cmNlcyA8PAovRm9udCA8PAov"
    "RjEgNSAwIFIKPj4KPj4KPj4KZW5kb2JqCjQgMCBvYmoKPDwKL0xlbmd0aCA0NQo+PgpzdHJlYW0K"
    "QlQKL0YxIDEyIFRmCjEwMCA3MDAgVGQoSGVsbG8sIFBERikgVGoKRVQKZW5kc3RyZWFtCmVuZG9i"
    "agoxIDAgb2JqCjw8Ci9UeXBlIC9Gb250Ci9TdWJ0eXBlIC9UeXBlMQovQmFzZUZvbnQgL0hlbHZl"
    "dGljYQo+PgplbmRvYmoKeHJlZgowIDYKMDAwMDAwMDAwIDY1NTM1IGYgCjAwMDAwMDAxMSAwMDAw"
    "MCBuIAowMDAwMDAwMTAwIDAwMDAwIG4gCjAwMDAwMDAxNjQgMDAwMDAgbiAKMDAwMDAwMDI1MCAw"
    "MDAwMCBuIAowMDAwMDAwMzM1IDAwMDAwIG4gCnRyYWlsZXIKPDwKL1NpemUgNgovUm9vdCAxIDAg"
    "Ugo+PgpzdGFydHhyZWYKNDI1CiUlRU9G"
)

_SAMPLE_PNG_BASE64 = (
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR4nGNgYAAAAAMAASsJTYQAAA"
    "AASUVORK5CYII="
)


class TestConvertToMarkdownTool(unittest.TestCase):
    """Validate conversion for multiple file types."""

    def setUp(self):
        self.tool = ConvertToMarkdownTool()

    @unittest.skipIf(docx is None, "python-docx not installed")
    def test_convert_docx(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            doc_path = Path(tmpdir) / "sample.docx"
            document = docx.Document()
            document.add_heading("Docx Title", level=1)
            document.add_paragraph("First paragraph.")
            document.add_paragraph("Second paragraph with more text.")
            table = document.add_table(rows=2, cols=2)
            table.rows[0].cells[0].text = "Header A"
            table.rows[0].cells[1].text = "Header B"
            table.rows[1].cells[0].text = "Value A"
            table.rows[1].cells[1].text = "Value B"
            document.save(str(doc_path))

            result = self.tool.execute(
                {"file_id": "docx_test", "source_path": str(doc_path), "file_type": "docx"}
            )
            markdown = result["raw_markdown"]
            self.assertIn("# Docx Title", markdown)
            self.assertIn("First paragraph.", markdown)
            self.assertIn("| Header A | Header B |", markdown)

    @unittest.skipIf(Presentation is None, "python-pptx not installed")
    def test_convert_pptx(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = Path(tmpdir) / "sample.pptx"
            presentation = Presentation()
            slide_layout = presentation.slide_layouts[1]
            slide = presentation.slides.add_slide(slide_layout)
            slide.shapes.title.text = "Slide Heading"
            slide.placeholders[1].text = "Bullet 1\nBullet 2"
            presentation.save(str(pptx_path))

            result = self.tool.execute(
                {"file_id": "pptx_test", "source_path": str(pptx_path), "file_type": "pptx"}
            )
            markdown = result["raw_markdown"]
            self.assertIn("## Slide Heading", markdown)
            self.assertIn("- Bullet 1", markdown)

    @unittest.skipIf(pdfplumber is None, "pdfplumber not installed")
    def test_convert_pdf(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            pdf_path = Path(tmpdir) / "sample.pdf"
            pdf_path.write_bytes(base64.b64decode(_SAMPLE_PDF_BASE64))

            result = self.tool.execute(
                {"file_id": "pdf_test", "source_path": str(pdf_path), "file_type": "pdf"}
            )
            markdown = result["raw_markdown"]
            self.assertIn("Hello, PDF", markdown)
            self.assertNotIn("## Page 1", markdown)


class TestExtractImagesTool(unittest.TestCase):
    """Validate image extraction for supported file types."""

    def setUp(self):
        self.tool = ExtractImagesTool()

    @unittest.skipIf(docx is None, "python-docx not installed")
    def test_extract_images_docx(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            image_path = Path(tmpdir) / "small.png"
            image_path.write_bytes(base64.b64decode(_SAMPLE_PNG_BASE64))

            doc_path = Path(tmpdir) / "image.docx"
            document = docx.Document()
            document.add_paragraph("With image below:")
            document.add_picture(str(image_path))
            document.save(str(doc_path))

            result = self.tool.execute(
                {"file_id": "docx_images", "source_path": str(doc_path), "file_type": "docx"}
            )
            self.assertGreater(len(result["images"]), 0)
            for image_meta in result["images"]:
                self.assertTrue(Path(image_meta["path"]).exists())

    @unittest.skipIf(Presentation is None, "python-pptx not installed")
    def test_extract_images_pptx(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            image_path = Path(tmpdir) / "small.png"
            image_path.write_bytes(base64.b64decode(_SAMPLE_PNG_BASE64))

            pptx_path = Path(tmpdir) / "image.pptx"
            presentation = Presentation()
            blank_layout = presentation.slide_layouts[5]
            slide = presentation.slides.add_slide(blank_layout)
            slide.shapes.add_picture(str(image_path), left=1000000, top=1000000, width=914400, height=914400)
            presentation.save(str(pptx_path))

            result = self.tool.execute(
                {"file_id": "pptx_images", "source_path": str(pptx_path), "file_type": "pptx"}
            )
            self.assertGreater(len(result["images"]), 0)
            for image_meta in result["images"]:
                self.assertTrue(Path(image_meta["path"]).exists())


if __name__ == "__main__":  # pragma: no cover
    unittest.main(verbosity=2)

